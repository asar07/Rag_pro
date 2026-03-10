import streamlit as st
import os
import tempfile
import math
import requests
import json

import pdfplumber
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup

st.set_page_config(page_title="DocChat", page_icon="📄", layout="wide")


# ---------------- API KEY ----------------

def get_api_key():
    try:
        return st.secrets["BYTEZ_API_KEY"]
    except Exception:
        return None


# ---------------- MODELS ----------------

MODELS = [
    "openai/gpt-4o",
    "openai/gpt-4o-mini",
]


# ---------------- CLEAN MODEL OUTPUT ----------------

def clean_output(raw):
    if isinstance(raw, str):
        try:
            parsed = json.loads(raw)
            if isinstance(parsed, dict) and "content" in parsed:
                return parsed["content"]
        except Exception:
            pass
        if "'content':" in raw:
            try:
                return raw.split("'content':")[1].split("'}")[0].strip(" '")
            except Exception:
                pass
        return raw
    if isinstance(raw, dict):
        if "content" in raw:
            return raw["content"]
        if "generated_text" in raw:
            return raw["generated_text"]
    if isinstance(raw, list):
        for item in raw:
            if isinstance(item, dict):
                if "message" in item:
                    return item["message"]["content"]
                if "generated_text" in item:
                    return item["generated_text"]
    return str(raw)


# ---------------- PDF ----------------

def extract_pdf(file_bytes):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        path = tmp.name
    pages = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append({"page": i + 1, "text": text.strip()})
        return pages
    finally:
        os.unlink(path)


# ---------------- DOCX ----------------

def extract_docx(file_bytes):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        tmp.write(file_bytes)
        path = tmp.name
    try:
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return [{"page": 1, "text": text}]
    finally:
        os.unlink(path)


# ---------------- PPTX ----------------

def extract_pptx(file_bytes):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(file_bytes)
        path = tmp.name
    slides = []
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            content = "\n".join(text)
            if content.strip():
                slides.append({"page": i + 1, "text": content})
        return slides
    finally:
        os.unlink(path)


# ---------------- TXT ----------------

def extract_txt(file_bytes):
    text = file_bytes.decode("utf-8", errors="ignore")
    return [{"page": 1, "text": text}]


# ---------------- URL ----------------

def extract_url(url):
    try:
        res = requests.get(url, timeout=10)
        soup = BeautifulSoup(res.text, "html.parser")
        for tag in soup(["script", "style"]):
            tag.extract()
        text = soup.get_text(separator=" ")
        return [{"page": 1, "text": text}]
    except Exception:
        return [{"page": 1, "text": ""}]


# ---------------- CHUNKING ----------------

def chunk_pages(pages, size=400, overlap=60):
    chunks = []
    for p in pages:
        words = p["text"].split()
        start = 0
        while start < len(words):
            end = min(start + size, len(words))
            chunks.append({
                "text": " ".join(words[start:end]),
                "page": p["page"],
                "id": len(chunks)
            })
            if end == len(words):
                break
            start += size - overlap
    return chunks


# ---------------- VECTOR INDEX ----------------

def build_index(chunks):
    vocab = {}
    for c in chunks:
        for w in c["text"].lower().split():
            if w not in vocab:
                vocab[w] = len(vocab)

    N = len(chunks)
    df = [0] * len(vocab)
    for c in chunks:
        for w in set(c["text"].lower().split()):
            if w in vocab:
                df[vocab[w]] += 1

    idf = [math.log((N + 1) / (d + 1)) + 1 for d in df]

    def vec(text):
        words = text.lower().split()
        tf = {}
        for w in words:
            tf[w] = tf.get(w, 0) + 1
        n = len(words) or 1
        v = [0.0] * len(vocab)
        for w, cnt in tf.items():
            if w in vocab:
                v[vocab[w]] = (cnt / n) * idf[vocab[w]]
        return v

    vecs = [vec(c["text"]) for c in chunks]
    return vocab, vecs, vec


# ---------------- COSINE ----------------

def cosine(a, b):
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = sum(x ** 2 for x in a) ** 0.5
    norm_b = sum(x ** 2 for x in b) ** 0.5
    return dot / (norm_a * norm_b + 1e-9)


# ---------------- RETRIEVE ----------------

def retrieve(query, chunks, vecs, vec_fn, k=4):
    qv = vec_fn(query)
    scored = sorted(
        enumerate(vecs),
        key=lambda x: cosine(qv, x[1]),
        reverse=True
    )
    return [chunks[i] for i, _ in scored[:k]]


# ---------------- MODEL CALL ----------------

def ask(question, context_chunks, history, api_key, model_id):
    context = "\n\n".join(
        f"(Page {c['page']}) {c['text']}"
        for c in context_chunks
    )

    system_prompt = f"""You are an expert document analyst and research assistant with deep expertise in extracting, synthesizing, and explaining information from all kinds of documents — reports, presentations, research papers, legal texts, manuals, and more.

## YOUR CORE MISSION
Provide thorough, insightful, and well-structured answers that go beyond surface-level retrieval. Think critically about the content, connect related ideas, and help the user truly understand what the document contains.

## RESPONSE QUALITY STANDARDS
- **Depth**: Give complete, detailed answers. Never truncate when the user deserves a full explanation.
- **Structure**: Use clear formatting — headers, bullet points, numbered lists, bold key terms — to make responses easy to scan.
- **Insight**: Don't just quote the document. Synthesize, explain implications, and highlight important patterns or themes.
- **Citations**: Always cite page/slide numbers inline like **(Page 3)** or **(Slide 7)** so users can verify.
- **Examples**: When the document contains examples, data, or figures that support your answer, include them.
- **Completeness**: If a question touches multiple parts of the document, address all relevant parts.

## RESPONSE FORMAT — adapt to the question type:
- **Factual questions** → Direct answer first, then supporting detail with citations.
- **Explanatory questions** → Break down concepts step by step with context from the document.
- **Summary requests** → Structured overview: key themes, main points, important details.
- **Comparison questions** → Use a table or side-by-side structure.
- **List/enumeration** → Numbered or bulleted lists with a brief explanation for each item.
- **Analysis questions** → Discuss significance, implications, or connections between ideas.

## STRICT RULES
1. Base ALL answers ONLY on the provided document context.
2. If information is absent from the document, clearly state: "The document does not contain information about this."
3. Never fabricate, invent, or assume facts not present in the context.
4. If a question is ambiguous, answer the most likely interpretation and note alternatives.
5. For follow-up questions, maintain context from the conversation history.
6. End longer responses with a **💡 Key Takeaway** — one sentence summarizing the most important point.

---
## DOCUMENT CONTEXT:
{context}
"""

    messages = []
    for m in history[-10:]:
        messages.append(m)
    messages.append({"role": "user", "content": question})

    try:
        response = requests.post(
            "https://api.bytez.com/models/v2/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json={
                "model": model_id,
                "messages": [{"role": "system", "content": system_prompt}] + messages,
                "max_tokens": 3000,
            },
            timeout=60,
        )

        if response.status_code != 200:
            return f"API error {response.status_code}: {response.text}"

        data = response.json()

        if "choices" in data:
            return data["choices"][0]["message"]["content"]

        return clean_output(data)

    except requests.exceptions.Timeout:
        return "Request timed out. Please try again."
    except Exception as e:
        return f"API error: {str(e)}"


# ---------------- SESSION STATE ----------------

for key, default in [
    ("chunks", []),
    ("vecs", []),
    ("vec_fn", None),
    ("history", []),
    ("doc_name", None),      # track which doc is loaded
    ("doc_loaded", False),   # flag so we don't reprocess on rerun
]:
    if key not in st.session_state:
        st.session_state[key] = default


# ================================================
# SIDEBAR
# ================================================

with st.sidebar:
    st.title("📄 DocChat")
    st.markdown("---")

    # API Key
    secret_key = get_api_key()
    if secret_key:
        api_key = secret_key
        st.success("API key loaded ✓", icon="🔑")
    else:
        api_key = st.text_input("Bytez API Key", type="password")

    st.markdown("---")

    # Upload
    st.subheader("📂 Load Document")
    uploaded = st.file_uploader(
        "Upload a file",
        type=["pdf", "docx", "pptx", "txt"],
        label_visibility="collapsed"
    )
    url = st.text_input("Or paste a URL", placeholder="https://...")

    st.markdown("---")

    # Settings
    st.subheader("⚙️ Settings")
    model_id = st.selectbox("Model", MODELS)
    chunk_size = st.slider("Chunk Size", 200, 800, 400)
    overlap = st.slider("Overlap", 0, 150, 60)
    top_k = st.slider("Top Chunks", 2, 8, 4)

    st.markdown("---")

    # Clear chat button
    if st.button("🗑️ Clear Chat", use_container_width=True):
        st.session_state.history = []
        st.rerun()

    # Doc status
    if st.session_state.doc_name:
        st.info(f"📎 {st.session_state.doc_name}")


# ================================================
# LOAD DOCUMENT — only process when a NEW file arrives
# ================================================

# Detect new upload by comparing file name to last loaded doc
if uploaded is not None:
    if uploaded.name != st.session_state.doc_name:
        file_bytes = uploaded.read()
        name = uploaded.name.lower()
        with st.spinner(f"Extracting {uploaded.name}..."):
            if name.endswith(".pdf"):
                pages = extract_pdf(file_bytes)
            elif name.endswith(".docx"):
                pages = extract_docx(file_bytes)
            elif name.endswith(".pptx"):
                pages = extract_pptx(file_bytes)
            elif name.endswith(".txt"):
                pages = extract_txt(file_bytes)
            else:
                pages = []

        if pages:
            chunks = chunk_pages(pages, chunk_size, overlap)
            if chunks:
                vocab, vecs, vec_fn = build_index(chunks)
                st.session_state.chunks = chunks
                st.session_state.vecs = vecs
                st.session_state.vec_fn = vec_fn
                st.session_state.history = []
                st.session_state.doc_name = uploaded.name
                st.session_state.doc_loaded = True
            else:
                st.warning("No text could be extracted from the document.")
        else:
            st.warning("Could not read the document.")

elif url and url != st.session_state.doc_name:
    with st.spinner("Fetching URL..."):
        pages = extract_url(url)
    if pages and pages[0]["text"].strip():
        chunks = chunk_pages(pages, chunk_size, overlap)
        if chunks:
            vocab, vecs, vec_fn = build_index(chunks)
            st.session_state.chunks = chunks
            st.session_state.vecs = vecs
            st.session_state.vec_fn = vec_fn
            st.session_state.history = []
            st.session_state.doc_name = url
            st.session_state.doc_loaded = True


# ================================================
# MAIN — CHAT AREA
# ================================================

st.title("💬 Chat with your Document")

if not st.session_state.chunks:
    st.info("👈 Upload a document or paste a URL in the sidebar to get started.")
else:
    # Render chat history
    for msg in st.session_state.history:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    # Chat input (always visible at bottom once doc is loaded)
    prompt = st.chat_input("Ask anything about the document...")

    if prompt:
        if not api_key:
            st.error("Please enter your Bytez API key in the sidebar.")
        else:
            st.session_state.history.append({"role": "user", "content": prompt})

            with st.chat_message("user"):
                st.markdown(prompt)

            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    srcs = retrieve(
                        prompt,
                        st.session_state.chunks,
                        st.session_state.vecs,
                        st.session_state.vec_fn,
                        top_k,
                    )
                    answer = ask(
                        prompt,
                        srcs,
                        st.session_state.history[:-1],
                        api_key,
                        model_id,
                    )
                st.markdown(answer)

            st.session_state.history.append({"role": "assistant", "content": answer})
            st.rerun()
        
